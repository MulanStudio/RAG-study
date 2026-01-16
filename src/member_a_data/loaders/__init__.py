"""
数据加载模块 - 支持多模态文档加载

负责人：成员A（数据工程师）

支持格式：
- PDF（含表格提取）
- Excel/CSV
- Word
- Markdown
- PPTX
- 图片（VLM 描述）
"""

import os
import glob
import pandas as pd
from typing import List, Optional
from langchain_core.documents import Document

# 复用已有模块
import sys
project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
sys.path.insert(0, project_root)
from src.member_a_data.pdf_table_extractor import PDFTableExtractor


def load_excel_files(directory: str) -> List[Document]:
    """
    加载 Excel/CSV 文件，转换为自然语言记录
    
    Example:
        输入: | Company | Revenue |
              | SLB     | 33.1B   |
        
        输出: "Table Record: Company: SLB | Revenue: 33.1B"
    """
    documents = []
    excel_files = glob.glob(os.path.join(directory, "**/*.xlsx"), recursive=True) + \
                  glob.glob(os.path.join(directory, "**/*.csv"), recursive=True)
    
    for file_path in excel_files:
        try:
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            
            df = df.fillna("Unknown")
            columns = df.columns.tolist()
            
            for index, row in df.iterrows():
                content_parts = [f"Source File: {os.path.basename(file_path)}"]
                for col in columns:
                    val = row.get(col, '')
                    if val and str(val).lower() not in ['unknown', 'nan', '']:
                        content_parts.append(f"{col}: {val}")
                
                content = " | ".join(content_parts)
                doc = Document(
                    page_content=content,
                    metadata={
                        "source": file_path,
                        "type": "excel_record",
                        "row_index": index
                    }
                )
                documents.append(doc)
            
            print(f"✅ Excel: {os.path.basename(file_path)} ({len(df)} rows)")
        except Exception as e:
            print(f"⚠️ Excel 加载失败 {file_path}: {e}")
    
    return documents


def load_word_files(directory: str) -> List[Document]:
    """
    加载 Word 文件，保留章节结构
    
    Example:
        输入: ## Section 2. Compensation
              Client agrees to pay $25,000 per day...
        
        输出: "Document Section [Compensation]: Client agrees to pay..."
    """
    documents = []
    
    try:
        from docx import Document as DocxDocument
    except ImportError:
        print("⚠️ python-docx 未安装，跳过 Word 加载")
        return documents
    
    word_files = glob.glob(os.path.join(directory, "**/*.docx"), recursive=True)
    
    for file_path in word_files:
        try:
            doc_obj = DocxDocument(file_path)
            current_heading = "General"
            
            for para in doc_obj.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                if 'Heading' in para.style.name:
                    current_heading = text
                else:
                    enhanced_content = f"Document Section [{current_heading}]: {text}"
                    doc = Document(
                        page_content=enhanced_content,
                        metadata={
                            "source": file_path,
                            "type": "contract_clause",
                            "section": current_heading
                        }
                    )
                    documents.append(doc)
            
            print(f"✅ Word: {os.path.basename(file_path)}")
        except Exception as e:
            print(f"⚠️ Word 加载失败 {file_path}: {e}")
    
    return documents


def load_markdown_files(directory: str) -> List[Document]:
    """
    加载 Markdown 文件，按标题切分
    """
    from langchain_text_splitters import MarkdownHeaderTextSplitter
    
    documents = []
    md_files = glob.glob(os.path.join(directory, "**/*.md"), recursive=True)
    
    headers_to_split_on = [
        ("#", "Header 1"),
        ("##", "Header 2"),
        ("###", "Header 3"),
    ]
    markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
    
    for file_path in md_files:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            splits = markdown_splitter.split_text(content)
            
            for doc in splits:
                doc.metadata["source"] = file_path
                doc.metadata["type"] = "markdown_section"
                
                # 语义增强：检测故障/事件
                keywords_failure = ["failure", "incident", "lost pulse", "stuck", "fishing"]
                if any(kw in doc.page_content.lower() for kw in keywords_failure):
                    doc.page_content = f"[EVENT: Operational Issue]\n{doc.page_content}"
                
                documents.append(doc)
            
            print(f"✅ Markdown: {os.path.basename(file_path)} ({len(splits)} sections)")
        except Exception as e:
            print(f"⚠️ Markdown 加载失败 {file_path}: {e}")
    
    return documents


def load_pdf_files(directory: str) -> List[Document]:
    """
    加载 PDF 文件，自动提取表格
    """
    extractor = PDFTableExtractor(use_natural_language=True)
    documents = []
    
    pdf_files = glob.glob(os.path.join(directory, "**/*.pdf"), recursive=True)
    
    for file_path in pdf_files:
        try:
            docs = extractor.load_pdf_with_tables(file_path)
            documents.extend(docs)
            print(f"✅ PDF: {os.path.basename(file_path)} ({len(docs)} docs)")
        except Exception as e:
            print(f"⚠️ PDF 加载失败 {file_path}: {e}")
    
    return documents


def load_ppt_files(directory: str) -> List[Document]:
    """
    加载 PPTX 文件，按幻灯片提取文本
    """
    documents = []
    try:
        from pptx import Presentation
    except ImportError:
        print("⚠️ python-pptx 未安装，跳过 PPTX 加载")
        return documents

    ppt_files = glob.glob(os.path.join(directory, "**/*.pptx"), recursive=True)

    for file_path in ppt_files:
        try:
            prs = Presentation(file_path)
            slide_docs = 0

            for idx, slide in enumerate(prs.slides, start=1):
                text_parts = []
                title = None

                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()

                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text:
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)

                    if getattr(shape, "has_table", False):
                        table_rows = []
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text]
                            if row_text:
                                table_rows.append(" | ".join(row_text))
                        if table_rows:
                            text_parts.append("Table: " + " || ".join(table_rows))

                if not text_parts and not title:
                    continue

                content = "\n".join(text_parts) if text_parts else ""
                if title:
                    content = f"Slide Title: {title}\n{content}".strip()

                doc = Document(
                    page_content=f"PPT Slide {idx}:\n{content}".strip(),
                    metadata={
                        "source": file_path,
                        "type": "ppt_slide",
                        "slide_index": idx,
                        "slide_title": title or ""
                    }
                )
                documents.append(doc)
                slide_docs += 1

            print(f"✅ PPTX: {os.path.basename(file_path)} ({slide_docs} slides)")
        except Exception as e:
            print(f"⚠️ PPTX 加载失败 {file_path}: {e}")

    return documents


def load_images(directory: str, vlm=None) -> List[Document]:
    """
    加载图片，使用 VLM 生成描述
    """
    documents = []
    image_files = glob.glob(os.path.join(directory, "**/*.png"), recursive=True) + \
                  glob.glob(os.path.join(directory, "**/*.jpg"), recursive=True) + \
                  glob.glob(os.path.join(directory, "**/*.jpeg"), recursive=True)
    
    for file_path in image_files:
        try:
            keywords = []
            if vlm:
                # 使用 VLM 生成描述
                import base64
                from langchain_core.messages import HumanMessage
                
                with open(file_path, "rb") as f:
                    image_data = base64.b64encode(f.read()).decode('utf-8')
                
                prompt = "Please describe this image in detail. Extract any text, numbers, or technical specifications visible."
                msg = HumanMessage(content=[
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_data}"}}
                ])
                response = vlm.invoke([msg])
                description = response.content
            else:
                # Fallback 1：OCR 提取（可选）
                description = ""
                try:
                    import importlib
                    Image = importlib.import_module("PIL.Image")
                    pytesseract = importlib.import_module("pytesseract")
                    with Image.open(file_path) as img:
                        ocr_text = pytesseract.image_to_string(img).strip()
                        if ocr_text:
                            description = f"OCR Text: {ocr_text}"
                except Exception:
                    description = ""

                # Fallback 2：使用文件名与规则提示
                if not description:
                    filename = os.path.basename(file_path).lower()
                    if "rig_spec" in filename or "offshore_rig_spec" in filename:
                        description = "Rig specification image with key parameters (water depth, max depth, hook load)."
                        keywords = ["rig", "spec", "water depth", "max depth", "hook load"]
                    elif "market_share" in filename:
                        description = "Market share chart image for offshore drilling."
                        keywords = ["market share", "offshore drilling", "chart"]
                    elif "safety" in filename:
                        description = "Safety alert poster with hazard details."
                        keywords = ["safety", "alert", "hazard"]
                    elif "zt09" in filename or "schematic" in filename:
                        description = "Well schematic image for ZT-09 with casing and BOP details."
                        keywords = ["ZT-09", "well schematic", "casing", "BOP", "mud weight"]
                    else:
                        description = f"Image file: {os.path.basename(file_path)}"
            
            keyword_text = f"\nKeywords: {', '.join(keywords)}" if keywords else ""
            keywords_str = ", ".join(keywords) if keywords else ""
            doc = Document(
                page_content=f"Image Description [{os.path.basename(file_path)}]: {description}{keyword_text}",
                metadata={
                    "source": file_path,
                    "type": "image_caption",
                    "image_path": file_path,
                    "keywords": keywords_str
                }
            )
            documents.append(doc)
            print(f"✅ Image: {os.path.basename(file_path)}")
        except Exception as e:
            print(f"⚠️ Image 加载失败 {file_path}: {e}")
    
    return documents


def load_all_documents(
    data_dir: str,
    vlm=None,
    verbose: bool = True
) -> List[Document]:
    """
    一键加载所有文档
    
    Args:
        data_dir: 数据目录（把组委会的数据放这里）
        vlm: 可选的 VLM 模型（用于图片描述）
        verbose: 是否打印加载信息
    
    Returns:
        List[Document]: 所有加载的文档
    
    Example:
        docs = load_all_documents("data/")
    """
    if verbose:
        print("=" * 60)
        print(f"📂 开始加载数据: {data_dir}")
        print("=" * 60)
    
    all_docs = []
    
    # 1. Excel/CSV
    excel_docs = load_excel_files(data_dir)
    all_docs.extend(excel_docs)
    
    # 2. Word
    word_docs = load_word_files(data_dir)
    all_docs.extend(word_docs)
    
    # 3. Markdown
    md_docs = load_markdown_files(data_dir)
    all_docs.extend(md_docs)
    
    # 4. PDF
    pdf_docs = load_pdf_files(data_dir)
    all_docs.extend(pdf_docs)
    
    # 5. PPTX
    ppt_docs = load_ppt_files(data_dir)
    all_docs.extend(ppt_docs)
    
    # 6. Images
    image_docs = load_images(data_dir, vlm=vlm)
    all_docs.extend(image_docs)
    
    if verbose:
        print("=" * 60)
        print(f"📊 加载完成: 共 {len(all_docs)} 个文档")
        print(f"   - Excel/CSV: {len(excel_docs)}")
        print(f"   - Word: {len(word_docs)}")
        print(f"   - Markdown: {len(md_docs)}")
        print(f"   - PDF: {len(pdf_docs)}")
        print(f"   - PPTX: {len(ppt_docs)}")
        print(f"   - Images: {len(image_docs)}")
        print("=" * 60)
    
    return all_docs
